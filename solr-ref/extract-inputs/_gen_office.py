#!/usr/bin/env python3
# Generates tiny, valid office documents + RTF for the issue #260 Solr captures.
# Also writes a malformed (truncated) variant of each binary format for the
# captured 500 envelope. Run once from this directory; outputs land here.
#
# These are INPUT fixtures, not ground truth -- the ground truth is whatever
# solr:9.10.1's Tika emits for them, captured into solr-ref/responses/.
import datetime
import os
import struct
import zipfile

from docx import Document
from docx.document import Document as _Doc
from pptx import Presentation
from openpyxl import Workbook
from odf.opendocument import OpenDocumentText, OpenDocumentSpreadsheet, OpenDocumentPresentation
from odf.style import Style, TextProperties, ParagraphProperties
from odf.text import P, H

CREATED = datetime.datetime(2024, 3, 14, 9, 26, 5, tzinfo=datetime.timezone.utc)

DOCX_TITLE = "Office Capture Title"
DOCX_AUTHOR = "Ada Example"


def write_docx():
    doc = Document()
    doc.core_properties.title = DOCX_TITLE
    doc.core_properties.author = DOCX_AUTHOR
    doc.core_properties.created = CREATED
    doc.core_properties.modified = CREATED
    doc.add_heading("Docx Heading", level=1)
    doc.add_paragraph("First DOCX paragraph.")
    doc.add_paragraph("Second DOCX paragraph with more text.")
    doc.save("sample.docx")


def write_pptx():
    prs = Presentation()
    prs.core_properties.title = DOCX_TITLE
    prs.core_properties.author = DOCX_AUTHOR
    prs.core_properties.created = CREATED
    prs.core_properties.modified = CREATED
    slide_layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide One Title"
    slide.placeholders[1].text = "Slide one bullet."
    body_layout = prs.slide_layouts[1]  # title + content
    slide2 = prs.slides.add_slide(body_layout)
    slide2.shapes.title.text = "Slide Two Title"
    slide2.placeholders[1].text = "Slide two content."
    prs.save("sample.pptx")


def write_xlsx():
    wb = Workbook()
    wb.properties.title = DOCX_TITLE
    wb.properties.creator = DOCX_AUTHOR
    wb.properties.created = CREATED
    wb.properties.modified = CREATED
    ws = wb.active
    ws.title = "Sheet One"
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "Alpha"
    ws["B2"] = 1
    ws["A3"] = "Beta"
    ws["B3"] = 2
    ws2 = wb.create_sheet("Sheet Two")
    ws2["A1"] = "Second sheet cell."
    wb.save("sample.xlsx")


def write_ods():
    from odf.table import Table, TableCell, TableRow
    doc = OpenDocumentSpreadsheet()
    doc.meta.title = DOCX_TITLE
    doc.meta.initial_creator = DOCX_AUTHOR
    doc.meta.creation_date = CREATED
    tbl = Table(name="Sheet One")
    cell = TableCell(valuetype="string")
    cell.addElement(P(text="ODS cell text."))
    row = TableRow()
    row.addElement(cell)
    tbl.addElement(row)
    doc.spreadsheet.addElement(tbl)
    doc.save("sample.ods")


def write_odt():
    doc = OpenDocumentText()
    doc.meta.title = DOCX_TITLE
    doc.meta.initial_creator = DOCX_AUTHOR
    doc.meta.creation_date = CREATED
    doc.text.addElement(H(outlinelevel="1", text="ODT Heading"))
    doc.text.addElement(P(text="First ODT paragraph."))
    doc.text.addElement(P(text="Second ODT paragraph."))
    doc.save("sample.odt")


def write_odp():
    from odf.draw import Page, Frame, TextBox
    doc = OpenDocumentPresentation()
    doc.meta.title = DOCX_TITLE
    doc.meta.initial_creator = DOCX_AUTHOR
    doc.meta.creation_date = CREATED
    page = Page(stylename="dp1", masterpagename="Default")
    frame = Frame(width="24cm", height="3cm", x="1cm", y="1cm")
    tb = TextBox()
    tb.addElement(P(text="ODP slide content."))
    frame.addElement(tb)
    page.addElement(frame)
    doc.presentation.addElement(page)
    doc.save("sample.odp")


def write_rtf():
    # Minimal hand-written RTF. ASCII control words, one bold run.
    # {\rtf1\ansi <title metadata is not in RTF core>; content.}
    rtf = (
        "{\\rtf1\\ansi\\deff0\n"
        "{\\fonttbl{\\f0 Times New Roman;}}\n"
        "\\f0\\fs24 First RTF paragraph.\\par\n"
        "{\\b Bold RTF run.}\\par\n"
        "Second RTF paragraph.\n"
        "}\n"
    )
    with open("sample.rtf", "w", encoding="ascii") as f:
        f.write(rtf)


def strip_thumbnail(path):
    """Remove docProps/thumbnail.jpeg so Tika emits no <div class=embedded>.

    python-docx/python-pptx embed a preview thumbnail; Tika surfaces it as an
    `<div class="embedded" id="/docProps/thumbnail.jpeg" />` at the end of
    the body. Stripping it keeps the fixture body to the document's actual
    text content.
    """
    import shutil
    tmp = path + ".tmp"
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            if item.filename == "docProps/thumbnail.jpeg":
                continue
            zout.writestr(item, zin.read(item.filename))
    shutil.move(tmp, path)


def truncate(src, dst, keep):
    """Write the first `keep` bytes of `src` to `dst` -- a corrupt archive."""
    with open(src, "rb") as f:
        data = f.read()
    with open(dst, "wb") as f:
        f.write(data[:keep])


def main():
    write_docx()
    write_pptx()
    write_xlsx()
    write_ods()
    write_odt()
    write_odp()
    write_rtf()
    for ext in ("docx", "pptx"):
        strip_thumbnail(f"sample.{ext}")
    # Malformed variants: truncating a zip mid-central-directory so Tika
    # throws (captured 500 envelope). RTF is plain text, so its malformed
    # form is a broken control-word stream Tika still chokes on.
    for ext in ("docx", "pptx", "xlsx", "ods", "odt", "odp"):
        truncate(f"sample.{ext}", f"broken.{ext}", 64)
    with open("broken.rtf", "w", encoding="ascii") as f:
        # A `\bin` keyword claiming ~10^10 bytes: Tika tries to skip that many
        # bytes and hits EOF, throwing EOFException -> captured 500 envelope.
        f.write("{\\rtf1\\ansi\\bin9999999999 ")

    # Zip-bomb fixture for the containment test
    # (`tests/office_extractor.rs::a_docx_shaped_zip_bomb_is_rejected_by_the_declared_ratio_guard`).
    # A DOCX-shaped package whose `word/document.xml` is 1 MiB of zeros: it
    # compresses to ~1 KiB, so the declared ratio (~1000:1) blows past the
    # 200:1 `zip_max_compression_ratio` guard at admission time, before a
    # byte is decompressed. The committed file is ~1 KiB on disk because
    # DEFLATE stores the compressed stream.
    import io
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0"?>\n<Types '
            'xmlns="http://schemas.openxmlformats.org/package/2006/content-types"/>',
        )
        z.writestr("word/document.xml", b"\0" * (1024 * 1024))
    with open("bomb.docx", "wb") as f:
        f.write(buf.getvalue())


if __name__ == "__main__":
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    main()
    for name in sorted(os.listdir(".")):
        if name.startswith(("_gen_office.py",)):
            continue
        if name.endswith((".docx", ".pptx", ".xlsx", ".ods", ".odt", ".odp", ".rtf")):
            print(f"{name}: {os.path.getsize(name)} bytes")
